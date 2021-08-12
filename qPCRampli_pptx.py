import os
import pandas as pd
from plotnine import *
import string
from scipy import stats
import numpy as np

xml_file_name = input("Path to xml file: ")       #S:\Experiments\E01_primer_screening\03_Infectious diseases\T015_20200805_FLUA_FLUB_gRNA\metadata\T015P01_metadata.xlsx
xlsx_file_name = input("Path to metadata xlsx file: ")      #'C:\\Users\\maria.rojec\\OneDrive - dnaNudge\\Documents\\Scripts\\Excel_metadata_example_T013P01.xlsx'
calcEff = input("Calculate PM efficiency? ")

#xml_file_name = 'S:\\Experiments\\E16_Respiratory_Virus_Panel\\A07_PCR_PMassay_validation\\T002_FLU_Inclusivity\\normalised_data\\calculated_data_T002P02.xml'
#xlsx_file_name = 'S:\\Experiments\\E16_Respiratory_Virus_Panel\\A07_PCR_PMassay_validation\\T002_FLU_Inclusivity\\metadata\\metadata_T002P02.xlsx'
#calcEff = 'n'

path_to_exp = xml_file_name.split('\\')[:5]
path_to_exp = '\\'.join(path_to_exp)

def num_list(start,end):
    return list(range(start,end+1))

def trim_columns(df):
    """
    Trim specific charcaters from ends of each value across all series in dataframe
    """
    trim_strings = lambda x: x[:-4] if isinstance(x, str) else x
    return df.applymap(trim_strings)

#Import xml file to extract normalised fluorescence measurements
import xml.etree.ElementTree as ET

tree = ET.parse(xml_file_name)

root = tree.getroot()

cq_list = []
for cycle in root.findall('.//{http://www.roche.ch/LC96AbsQuantCalculatedDataModel}absQuantDataSource'):
    cell_id = cycle.find('{http://www.roche.ch/LC96AbsQuantCalculatedDataModel}graphId').text
    if cycle.find('{http://www.roche.ch/LC96AbsQuantCalculatedDataModel}call').text == 'Positive':
        cq = cycle.find('{http://www.roche.ch/LC96AbsQuantCalculatedDataModel}cq').text
        cq_list.append([cell_id,float(cq)])

cq_df = pd.DataFrame(cq_list, columns=['Cell_id', 'Cq'])

norm_fluor = []
for cycle in root.findall('.//{http://www.roche.ch/LC96CalculatedDataSchema}factGraph'):
    cell_id = cycle.get('id')
    for num in cycle.findall('.//{http://www.roche.ch/LC96CalculatedDataSchema}fluor'):
        fluor = num.text
        norm_fluor.append([cell_id,float(fluor)])

cycle_number = int(len(norm_fluor)/96)

wells = []
#i=0
for a in range(8):
    row = string.ascii_uppercase[a]
    for m in range(1,13):
        cell = row+str(m)
        #cell_id = cells_id[i]
        #i += 1
        for i in range(cycle_number):
            wells.append(cell)

normFluor_df = pd.DataFrame(norm_fluor, columns=['Cell_id', 'normFluorescence'])
normFluor_df['Cells'] = wells

number_list = list(range(0,len(norm_fluor),cycle_number))
data = []
for i in range(96):
    start = number_list[i]
    end = number_list[i]+cycle_number
    cycle_list = list(normFluor_df['normFluorescence'][start:end])
    #cycle_list = norm_fluor[(number_list[i][1]):(number_list[i+cycle_number][1])]
    for l in range(cycle_number):
        cycle = l+1
        data.append([cycle,float(cycle_list[l])])
    #i+=1

cycle_df = pd.DataFrame(data, columns=['Cycle', 'normFluorescence'])

df_norm = pd.merge(normFluor_df, cycle_df)

df_norm = pd.merge(df_norm, cq_df)


##### Read metadata information
xls = pd.ExcelFile(xlsx_file_name)

#sheets = xls.sheet_names
res = len(xls.sheet_names)

# Make df machine readable by transforming it into long format, renaming and creating new cols
def machine_readable_df(df,string):
    df = df.rename(columns={"Unnamed: 0": "Rows"})
    df = pd.melt(df, id_vars =['Rows'], value_vars =num_list(1,12),var_name ='Columns', value_name =string.lstrip('df_')) 
    df['Cells'] = df["Rows"] + df["Columns"].astype(str)#
    return df

#Loop through all sheets in the excel file and append the metadata info to the data df
final_df = df_norm
for i in range(res):
    key_name = 'df_'+str(xls.sheet_names[i]) 
    df = pd.read_excel(xls, i)
    mr_df = machine_readable_df(df,key_name)
    final_df = pd.merge(final_df, mr_df)

#-----------------------------------------------------------------------------------------------------------------------
# Add images to presentation

from pptx import Presentation
from pptx.util import Inches
import os
 
prs = Presentation()

# front page
#-----------------------------------------------------------------------------------------------------------------------
layout0 = prs.slide_layouts[0]
slide = prs.slides.add_slide(layout0)
title=slide.shapes.title
subtitle=slide.placeholders[1]

title.text="Title" # title
subtitle.text="Maria Rojec" # subtitle

# PM slide
#-----------------------------------------------------------------------------------------------------------------------

def add_3image_slide(img_path,float1,float2,img_path2,float3,float4,img_path3,float5,float6):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    title=slide.shapes.title
    #title.text=string
    left1=float1
    top1=float2
    left2=float3
    top2=float4
    left3=float5
    top3=float6
    slide.shapes.add_picture(img_path,left1,top1)
    slide.shapes.add_picture(img_path2,left2,top2)
    slide.shapes.add_picture(img_path3,left3,top3)

def add_2image_slide(img_path,float1,float2,img_path2,float3,float4):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    title=slide.shapes.title
    #title.text=string
    left1=float1
    top1=float2
    left2=float3
    top2=float4
    slide.shapes.add_picture(img_path,left1,top1)
    slide.shapes.add_picture(img_path2,left2,top2)



#Plot amplification curves based on TARGET GENE.
#Also I may need to normalise raw data by subtracting baseline fluorescence, etc. 

#final_df = final_df.dropna()
final_df["Template_conc_log10"] = np.log10(final_df["Template_conc"])
final_df["Template_conc"] = final_df["Template_conc"].astype('category')

exp_plate = xlsx_file_name.split('\\')  
plate = exp_plate[len(exp_plate)-1].rstrip('.xlsx').lstrip('metadata_')

# Program to find most frequent  
# element in a list 
def most_frequent(List): 
    return max(set(List), key = List.count) 


for exp in final_df['Experiment'].unique():
    subset_df = final_df[final_df['Experiment'] == exp]
    sample = np.array2string(subset_df['Target'].unique()).strip("['']")
    template = most_frequent(list(subset_df['Template']))
    p=(ggplot(subset_df, aes(x='Cycle', y='normFluorescence',group='Cells',colour='factor(Template)')) +
            geom_line(size=1, alpha=0.5) + #scale_linetype_manual(values = ['-','--','-.', ]) +            
            scale_x_continuous(breaks=list(range(0,cycle_number+1,3))) + theme_bw() +
            labs(title=sample+' ('+plate+')', colour ="Template concentrations", linetype='Template') + theme(plot_title=element_text(size = 20, face = "bold")) + 
            ylab("Fluorescence"))
    image1_file_name = path_to_exp+'\\'+'AmplificationCurve_'+str(sample)+'_'+str(template)+'_'+plate+'.png'
    ggsave(plot = p, filename = image1_file_name, width=5, height=3)
    q=(ggplot(subset_df, aes(x='Cycle', y='normFluorescence',group='Cells',colour='factor(Template)')) +
            geom_line(size=1, alpha=0.5) + #scale_linetype_manual(values = ['-','--','-.']) +         
            scale_x_continuous(breaks=list(range(0,cycle_number+1,3))) + scale_y_log10(limits = [0.01, 2]) + theme_bw() +
            labs(title=sample+' ('+plate+')', colour ="Template concentrations", linetype='Template') + theme(plot_title=element_text(size = 20, face = "bold")) + 
            ylab("Fluorescence (log10)") + geom_hline(yintercept = 0.05,color='blue',size=0.5,linetype="dashed"))
    image2_file_name = path_to_exp+'\\'+'Log_AmplificationCurve_'+str(sample)+'_'+str(template)+'_'+plate+'.png'
    ggsave(plot = q, filename = image2_file_name, width=5, height=3)
    
    for sample in subset_df['Target'].unique():
        sample_df = subset_df[(subset_df['Target'] == sample) & (subset_df['Template'] == template)]
        if calcEff == 'Y' or calcEff == 'y' or calcEff == 'yes' or calcEff == 'Yes':
            sample_df = sample_df[sample_df['Cq'] > 15]
            slope, intercept, r_value, p_value, std_err = stats.linregress(sample_df['Template_conc_log10'],sample_df['Cq'])
            sample_df['fit']=sample_df.Template_conc_log10*slope+intercept
            efficiency = round(10**(-1/slope),2)
            #format text 
            txt= 'y = {:4.2e} x + {:4.2E};\n   R^2= {:2.2f}'.format(slope, intercept, r_value*r_value)
            #create plot. The 'factor' is a nice trick to force a discrete color scale
            plot=(ggplot(data=sample_df, mapping= aes('Template_conc_log10','Cq'))
                    + geom_point(aes()) + scale_y_continuous(breaks = list(range(15, 42,2)),limits = [18, 42])
                    + xlab("Template copy number (log10)")+ ylab(r'Cq')
                    + geom_line(aes(x='Template_conc_log10', y='fit'), color='black')
                    + theme_bw() + theme(title = element_text(color = "black", size = 10))
                    + labs(title=txt+'\n Efficiency: '+str(efficiency)))
            image3_file_name = path_to_exp+'\\'+'Efficiency_'+str(sample)+'_'+str(template)+'_'+plate+'.png'
            ggsave(plot = plot, filename = image3_file_name, width=2, height=2)
            add_3image_slide(image1_file_name,Inches(0),Inches(0.5),image2_file_name,Inches(0),Inches(4),image3_file_name,Inches(7.2),Inches(2.5))
        else:
            add_2image_slide(image1_file_name,Inches(0),Inches(0.5),image2_file_name,Inches(0),Inches(4))


prs.save(path_to_exp+'\\'+plate+".pptx") # saving file








